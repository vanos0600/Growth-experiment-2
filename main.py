"""
Revenue Optimization Engine
============================
A/B Test Intelligence Platform for Growth Teams.
Converts experiment data into automated revenue decisions.
"""

import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import json
from datetime import datetime

# ── PAGE CONFIG ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Revenue Optimization Engine",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ── GLOBAL CSS ────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600&family=JetBrains+Mono:wght@400;500&display=swap');

html, body, [class*="css"], .stApp {
    font-family: 'Inter', sans-serif !important;
    background-color: #080a0f !important;
    color: #dde1e7 !important;
}

.section-head {
    font-size: 10px; font-weight: 600; letter-spacing: .14em;
    text-transform: uppercase; color: #2d3650;
    font-family: 'JetBrains Mono', monospace;
    border-top: 1px solid #141824; padding-top: 20px; margin: 24px 0 12px;
}

.page-title {
    font-size: 26px; font-weight: 600;
    font-family: 'JetBrains Mono', monospace;
    color: #f0f4fa; letter-spacing: -.02em; margin-bottom: 2px;
}
.page-sub {
    font-size: 12px; color: #2d3650;
    font-family: 'JetBrains Mono', monospace; margin-bottom: 32px;
}

.hero-card {
    background: #0f1220; border: 0.5px solid #1c2035;
    border-radius: 14px; padding: 20px 18px;
    position: relative; overflow: hidden;
}
.hero-card::before {
    content: ''; position: absolute; top: 0; left: 0; right: 0; height: 3px;
    border-radius: 14px 14px 0 0;
}
.hero-card.green::before  { background: linear-gradient(90deg,#10b981,#34d399); }
.hero-card.blue::before   { background: linear-gradient(90deg,#3b82f6,#60a5fa); }
.hero-card.amber::before  { background: linear-gradient(90deg,#f59e0b,#fbbf24); }
.hero-card.red::before    { background: linear-gradient(90deg,#ef4444,#f87171); }
.hero-card.purple::before { background: linear-gradient(90deg,#8b5cf6,#a78bfa); }
.hero-card.slate::before  { background: linear-gradient(90deg,#475569,#64748b); }

.hero-label {
    font-size: 10px; font-weight: 600; letter-spacing: .1em; text-transform: uppercase;
    color: #3d4a6b; font-family: 'JetBrains Mono', monospace; margin-bottom: 10px;
}
.hero-value {
    font-size: 28px; font-weight: 600; line-height: 1.1;
    margin-bottom: 5px; font-family: 'JetBrains Mono', monospace;
}
.hero-value.green  { color: #10b981; }
.hero-value.blue   { color: #60a5fa; }
.hero-value.amber  { color: #f59e0b; }
.hero-value.red    { color: #ef4444; }
.hero-value.purple { color: #a78bfa; }
.hero-value.white  { color: #f0f4fa; }
.hero-sub { font-size: 11px; color: #2d3650; line-height: 1.5; }

.action-banner {
    border-radius: 12px; padding: 18px 22px; margin-bottom: 8px;
    display: flex; align-items: flex-start; gap: 16px;
}
.action-banner.critical { background: #120404; border: 1px solid #7f1d1d; }
.action-banner.warning  { background: #130e00; border: 1px solid #78350f; }
.action-banner.success  { background: #030f09; border: 1px solid #064e3b; }
.action-banner.info     { background: #030d1f; border: 1px solid #1e3a5f; }

.action-icon   { font-size: 22px; flex-shrink: 0; line-height: 1; margin-top: 1px; }
.action-title  { font-size: 13px; font-weight: 600; font-family: 'JetBrains Mono', monospace; margin-bottom: 3px; }
.action-title.critical { color: #f87171; }
.action-title.warning  { color: #fbbf24; }
.action-title.success  { color: #34d399; }
.action-title.info     { color: #60a5fa; }
.action-body   { font-size: 12px; color: #4b5780; line-height: 1.65; }
.action-body b { color: #8892a4; font-weight: 500; }

.insight-card {
    background: #0f1220; border: 0.5px solid #1c2035;
    border-radius: 12px; padding: 18px 20px; margin-bottom: 10px;
    font-size: 13px; color: #4b5780; line-height: 1.75;
}
.insight-card b { color: #dde1e7; font-weight: 500; }

.hyp-card {
    background: #0f1220; border: 0.5px solid #1c2035;
    border-left: 3px solid #8b5cf6;
    border-radius: 0 12px 12px 0;
    padding: 14px 18px; margin-bottom: 10px;
}
.hyp-num  { font-size: 10px; font-family: 'JetBrains Mono', monospace; color: #8b5cf6; font-weight: 600; margin-bottom: 4px; letter-spacing:.08em; }
.hyp-body { font-size: 13px; color: #8892a4; line-height: 1.6; }
.hyp-body b { color: #dde1e7; font-weight: 500; }

.var-row {
    background: #0f1220; border: 0.5px solid #1c2035;
    border-radius: 10px; padding: 14px 16px; margin-bottom: 8px;
}
.var-row.winner { border-color: #10b981; border-width: 1px; }
.var-name { font-size: 14px; font-weight: 500; color: #dde1e7; margin-bottom: 2px; }
.var-copy { font-size: 11px; color: #2d3650; font-style: italic; margin-bottom: 10px; }
.var-stats { display: flex; gap: 20px; flex-wrap: wrap; }
.vs-label { font-size: 10px; color: #2d3650; text-transform: uppercase; letter-spacing:.06em; font-family: 'JetBrains Mono', monospace; }
.vs-value { font-size: 13px; font-weight: 500; color: #dde1e7; font-family: 'JetBrains Mono', monospace; }
.vs-value.miss { color: #ef4444; }
.vs-value.win  { color: #10b981; }
.rpi-bar-track { background: #1c2035; height: 4px; border-radius: 2px; margin-top: 6px; }

.stDownloadButton > button {
    background: #10b981 !important; color: #022c22 !important;
    border: none !important; border-radius: 10px !important;
    font-weight: 700 !important; font-size: 14px !important;
    height: 46px !important; width: 100% !important;
}

.stTabs [data-baseweb="tab-list"] {
    background: #0f1220; border-bottom: 1px solid #1c2035; gap: 4px;
}
.stTabs [data-baseweb="tab"] {
    color: #3d4a6b !important; font-size: 13px !important; font-weight: 500 !important;
    padding: 8px 16px !important;
}
.stTabs [aria-selected="true"] {
    color: #dde1e7 !important; border-bottom: 2px solid #10b981 !important;
}
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# DEFAULT DATA
# ══════════════════════════════════════════════════════════════════════════════

DEFAULT_DF = pd.DataFrame([
    {"Variant": "A", "Copy": "Upgrade now for full protection",
     "Impressions": 120_000, "CTR_%": 0.50, "Conversion_%": 0.40,
     "Revenue_$": 2_500.0, "Notes": "Control"},
    {"Variant": "B", "Copy": "Stay safe online – unlock premium features",
     "Impressions": 110_000, "CTR_%": 0.45, "Conversion_%": 0.50,
     "Revenue_$": 3_800.0, "Notes": ""},
    {"Variant": "C", "Copy": "Your device may be at risk – fix it now",
     "Impressions": 120_000, "CTR_%": 0.0,  "Conversion_%": 1.20,
     "Revenue_$": 5_100.0, "Notes": "too scary (user feedback); CTR tracking failed"},
])

PLOTLY_DARK = dict(
    plot_bgcolor="#0f1220", paper_bgcolor="#0f1220",
    font_color="#8892a4", font_family="JetBrains Mono",
)


# ══════════════════════════════════════════════════════════════════════════════
# ANALYSIS ENGINE
# ══════════════════════════════════════════════════════════════════════════════

def enrich(df: pd.DataFrame) -> pd.DataFrame:
    d = df.copy()
    d["ctr_missing"] = d["CTR_%"] == 0.0
    d["RPI"]         = d["Revenue_$"] / d["Impressions"]
    max_rpi          = d["RPI"].max()
    d["RPI_pct"]     = d["RPI"] / max_rpi
    d["is_winner"]   = d["Revenue_$"] == d["Revenue_$"].max()
    notes_lower      = d["Notes"].str.lower()
    d["brand_risk"]  = notes_lower.apply(
        lambda n: "high"   if any(w in n for w in ["unsubscribe","toxic","angry","complaint"]) else
                  "medium" if any(w in n for w in ["scary","fear","alarming","negative","aggressive","too"]) else
                  "low"
    )
    return d


def revenue_guard(enriched: pd.DataFrame) -> dict:
    missing_ctr      = enriched[enriched["ctr_missing"]]
    brand_alerts     = enriched[enriched["brand_risk"] != "low"]
    total_imp        = enriched["Impressions"].sum()
    tracked_imp      = enriched[~enriched["ctr_missing"]]["Impressions"].sum()
    insight_gap_pct  = (1 - tracked_imp / total_imp) * 100 if total_imp else 0
    return {
        "missing_ctr_variants": missing_ctr["Variant"].tolist(),
        "brand_alert_variants": brand_alerts[["Variant","brand_risk"]].to_dict("records"),
        "insight_gap_pct":      insight_gap_pct,
        "tracking_ok":          len(missing_ctr) == 0,
    }


def opportunity_cost(enriched: pd.DataFrame, duration_days: int = 10) -> dict:
    winner             = enriched[enriched["is_winner"]].iloc[0]
    losers             = enriched[~enriched["is_winner"]]
    winner_rpi         = winner["RPI"]
    total_imp          = enriched["Impressions"].sum()
    loser_imp          = losers["Impressions"].sum()
    actual_revenue     = enriched["Revenue_$"].sum()
    potential_revenue  = winner_rpi * total_imp
    lost_revenue       = potential_revenue - actual_revenue
    daily_lost         = lost_revenue / duration_days
    if_scaled_monthly  = winner_rpi * (total_imp / duration_days) * 30
    baseline_monthly   = actual_revenue / duration_days * 30
    return {
        "winner_variant":       winner["Variant"],
        "winner_rpi":           winner_rpi,
        "winner_revenue":       winner["Revenue_$"],
        "winner_copy":          winner["Copy"],
        "actual_revenue":       actual_revenue,
        "potential_revenue":    potential_revenue,
        "lost_revenue":         lost_revenue,
        "daily_lost":           daily_lost,
        "monthly_lost":         daily_lost * 30,
        "annual_lost":          daily_lost * 365,
        "if_scaled_monthly":    if_scaled_monthly,
        "baseline_monthly":     baseline_monthly,
        "incremental_monthly":  if_scaled_monthly - baseline_monthly,
        "loser_imp":            loser_imp,
        "total_imp":            total_imp,
    }


def score_variant(row) -> float:
    penalty = {"low": 0.0, "medium": 0.15, "high": 0.40}[row["brand_risk"]]
    return round(row["RPI_pct"] * (1 - penalty) * 100, 1)


def decide_verdict(winner_row) -> dict:
    risk = winner_row["brand_risk"]
    if risk == "high":
        return {"verdict": "KILL",    "css": "critical", "icon": "🛑",
                "body": "Critical brand risk. Revenue gain does not justify reputational exposure. Revert to control immediately."}
    if risk == "medium":
        return {"verdict": "ITERATE", "css": "warning",  "icon": "🔁",
                "body": "Winner delivers +revenue but carries medium brand risk. Launch a tone-softened iteration before full rollout."}
    return         {"verdict": "SCALE",   "css": "success",  "icon": "🚀",
                "body": "No brand risk signals. Proceed with phased rollout: 20% → 50% → 100% over 9 days."}


# ══════════════════════════════════════════════════════════════════════════════
# CHARTS
# ══════════════════════════════════════════════════════════════════════════════

def chart_rpi(enriched):
    colors = ["#10b981" if w else "#1c2e4a" for w in enriched["is_winner"]]
    fig = go.Figure(go.Bar(
        x=enriched["Variant"], y=enriched["RPI"],
        marker_color=colors,
        text=[f"${v:.4f}" for v in enriched["RPI"]],
        textposition="outside",
        textfont=dict(size=11, family="JetBrains Mono"),
    ))
    fig.update_layout(**PLOTLY_DARK,
        title=dict(text="Revenue Per Impression — primary normalised metric",
                   font=dict(size=12, color="#4b5780"), x=0),
        height=260, margin=dict(l=0,r=0,t=40,b=0),
        xaxis=dict(showgrid=False, color="#3d4a6b"),
        yaxis=dict(showgrid=True, gridcolor="#141824", color="#3d4a6b", tickformat="$.4f"),
        showlegend=False,
    )
    return fig


def chart_waterfall(oc):
    fig = go.Figure(go.Waterfall(
        orientation="v",
        measure=["absolute", "relative", "total"],
        x=["Actual revenue\n(all variants)", "Lost to losing\nvariants", "Potential\n(winner only)"],
        y=[oc["actual_revenue"], oc["lost_revenue"], 0],
        totals={"marker": {"color": "#10b981"}},
        increasing={"marker": {"color": "#3b82f6"}},
        decreasing={"marker": {"color": "#ef4444"}},
        connector={"line": {"color": "#1c2035", "width": 1}},
        text=[f"${oc['actual_revenue']:,.0f}", f"${oc['lost_revenue']:,.0f}", f"${oc['potential_revenue']:,.0f}"],
        textfont=dict(size=11, family="JetBrains Mono"),
    ))
    fig.update_layout(**PLOTLY_DARK,
        title=dict(text="Revenue waterfall — money left on the table",
                   font=dict(size=12, color="#4b5780"), x=0),
        height=280, margin=dict(l=0,r=0,t=40,b=0),
        xaxis=dict(showgrid=False, color="#3d4a6b"),
        yaxis=dict(showgrid=True, gridcolor="#141824", color="#3d4a6b",
                   tickprefix="$", tickformat=",.0f"),
        showlegend=False,
    )
    return fig


def chart_funnel(enriched):
    fig = go.Figure()
    palette = {"A": "#3b82f6", "B": "#f59e0b", "C": "#10b981"}
    for _, row in enriched.iterrows():
        fig.add_trace(go.Bar(
            name=f"Variant {row['Variant']}",
            x=["CTR %", "Conversion %"],
            y=[row["CTR_%"] if not row["ctr_missing"] else 0, row["Conversion_%"]],
            marker_color=palette.get(row["Variant"], "#8892a4"),
        ))
    fig.update_layout(**PLOTLY_DARK,
        title=dict(text="Funnel comparison — CTR vs Conversion by variant",
                   font=dict(size=12, color="#4b5780"), x=0),
        height=270, margin=dict(l=0,r=0,t=40,b=0), barmode="group",
        xaxis=dict(showgrid=False, color="#3d4a6b"),
        yaxis=dict(showgrid=True, gridcolor="#141824", color="#3d4a6b", ticksuffix="%"),
        legend=dict(font=dict(color="#4b5780", size=11), bgcolor="rgba(0,0,0,0)"),
    )
    return fig


def chart_projection(oc):
    months   = list(range(1, 13))
    baseline = [oc["baseline_monthly"] * m for m in months]
    scaled   = [oc["if_scaled_monthly"] * m for m in months]
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=months, y=baseline, name="Current mix",
        line=dict(color="#3b82f6", width=2, dash="dot")))
    fig.add_trace(go.Scatter(x=months, y=scaled, name="Winner scaled",
        line=dict(color="#10b981", width=2),
        fill="tonexty", fillcolor="rgba(16,185,129,0.08)"))
    fig.update_layout(**PLOTLY_DARK,
        title=dict(text="12-month cumulative revenue — current mix vs winner scaled",
                   font=dict(size=12, color="#4b5780"), x=0),
        height=280, margin=dict(l=0,r=0,t=40,b=0),
        xaxis=dict(showgrid=False, color="#3d4a6b", title="Month"),
        yaxis=dict(showgrid=True, gridcolor="#141824", color="#3d4a6b",
                   tickprefix="$", tickformat=",.0f"),
        legend=dict(font=dict(color="#4b5780", size=11), bgcolor="rgba(0,0,0,0)"),
    )
    return fig


# ══════════════════════════════════════════════════════════════════════════════
# AI MODULE
# ══════════════════════════════════════════════════════════════════════════════

def build_prompt(enriched, oc, guard) -> str:
    winner    = enriched[enriched["is_winner"]].iloc[0]
    qual_rows = enriched[enriched["Notes"].str.strip() != ""][["Variant","Notes"]].to_dict("records")
    return f"""You are a Senior Growth Strategist and CRO expert.
Analyse the following A/B test results and generate actionable output.

EXPERIMENT DATA:
{enriched[['Variant','Copy','Impressions','CTR_%','Conversion_%','Revenue_$','brand_risk','RPI']].to_string(index=False)}

WINNER: Variant {winner['Variant']} — "{winner['Copy']}"
RPI: ${winner['RPI']:.4f} | Conversion: {winner['Conversion_%']}% | Revenue: ${winner['Revenue_$']:,.0f}

REVENUE OPPORTUNITY:
- Lost revenue this period: ${oc['lost_revenue']:,.0f}
- Incremental monthly if winner scaled: ${oc['incremental_monthly']:,.0f}
- Annual revenue at stake: ${oc['annual_lost']:,.0f}

QUALITATIVE SIGNALS:
{json.dumps(qual_rows, indent=2)}

TRACKING ISSUES:
{"CTR missing for: " + str(guard['missing_ctr_variants']) if not guard['tracking_ok'] else "None"}

Respond ONLY with valid JSON. No markdown, no fences, no extra text:
{{
  "executive_summary": "2-3 sentence management summary of the experiment outcome and key trade-off",
  "verdict": "SCALE | ITERATE | KILL",
  "verdict_rationale": "1 sentence explaining the verdict",
  "hypotheses": [
    {{
      "id": "H1",
      "copy": "Proposed copy ≤12 words",
      "rationale": "Why this preserves urgency while softening alarm",
      "predicted_outcome": "Expected metric improvement",
      "risk_level": "Low | Medium | High"
    }},
    {{"id": "H2", "copy": "...", "rationale": "...", "predicted_outcome": "...", "risk_level": "..."}},
    {{"id": "H3", "copy": "...", "rationale": "...", "predicted_outcome": "...", "risk_level": "..."}}
  ],
  "priority_actions": ["action 1", "action 2", "action 3"],
  "risk_mitigation": "Strategy to protect brand equity while scaling the winner"
}}

HYPOTHESIS RULES:
- All 3 must address Variant {winner['Variant']}'s feedback: "{winner['Notes']}"
- Maintain URGENCY, soften FEAR/ALARM tone
- Each copy ≤12 words, meaningfully distinct psychological triggers
"""


def make_rule_based_ai(enriched, oc, guard) -> dict:
    """Generate a structured report with pure Python logic — no API key needed."""
    winner  = enriched[enriched["is_winner"]].iloc[0]
    v       = decide_verdict(winner)
    lift    = (winner["Revenue_$"] - enriched.iloc[0]["Revenue_$"]) / enriched.iloc[0]["Revenue_$"] * 100
    risk    = winner["brand_risk"]

    summary = (
        f"Variant {winner['Variant']} is the clear financial winner with ${winner['Revenue_$']:,.0f} revenue, "
        f"a {winner['Conversion_%']}% conversion rate, and an RPI of ${winner['RPI']:.4f} — "
        f"a +{lift:.1f}% lift over the control. "
        f"The key trade-off is a {risk} brand risk signal from user feedback that must be resolved before scale. "
        f"Annual revenue at stake: ${oc['annual_lost']:,.0f}."
    )

    notes = winner["Notes"].lower()
    if any(w in notes for w in ["scary", "fear", "alarming", "too", "negative", "aggressive"]):
        hyps = [
            {"id": "H1",
             "copy": "Your device deserves better protection — upgrade now",
             "rationale": "Replaces fear framing with aspirational ownership. Keeps urgency ('now') without threat language.",
             "predicted_outcome": "Maintains ~1.0–1.1% conversion while reducing negative feedback by ~60–70%.",
             "risk_level": "Low"},
            {"id": "H2",
             "copy": "Most devices in your area are already protected — join them",
             "rationale": "Social proof + FOMO replaces direct fear. Urgency is implicit, not alarming.",
             "predicted_outcome": "0.9–1.1% conversion; social proof typically lifts CTR 10–20% over threat copy.",
             "risk_level": "Low"},
            {"id": "H3",
             "copy": "One tap to protect your device — takes 30 seconds",
             "rationale": "Urgency comes from ease, not danger. Removes 'risk' language entirely.",
             "predicted_outcome": "Slightly lower peak conversion (0.8–1.0%) but significantly lower brand-risk exposure.",
             "risk_level": "Low"},
        ]
    else:
        hyps = [
            {"id": "H1",
             "copy": "Upgrade today and unlock every feature — limited offer",
             "rationale": "Adds scarcity signal to test urgency amplification without fear language.",
             "predicted_outcome": "+5–10% conversion lift via scarcity framing.",
             "risk_level": "Low"},
            {"id": "H2",
             "copy": "Join 2 million users who upgraded — see what you're missing",
             "rationale": "Social proof reframe: community belonging replaces feature-listing.",
             "predicted_outcome": "Similar conversion, higher CTR due to social validation.",
             "risk_level": "Low"},
            {"id": "H3",
             "copy": "Your free trial ends soon — lock in full access now",
             "rationale": "Time-based urgency without fear. Tests deadline framing vs threat framing.",
             "predicted_outcome": "1.0–1.3% conversion depending on segment.",
             "risk_level": "Medium"},
        ]

    return {
        "executive_summary":  summary,
        "verdict":            v["verdict"],
        "verdict_rationale":  v["body"],
        "hypotheses":         hyps,
        "priority_actions": [
            f"Fix CTR tracking on Variant {winner['Variant']} before next iteration.",
            f"Launch H1 as Variant {winner['Variant']}.2 against the current winner within 7 days.",
            "Geo-validate the winner outside the current market before global rollout.",
        ],
        "risk_mitigation": (
            f"Run Variant {winner['Variant']}.2 (H1 copy) at 10% traffic alongside the winner. "
            "Monitor support tickets and NPS. Scale H1 when negative feedback drops below 2% of sessions."
        ),
    }


def ai_component(prompt: str) -> None:
    """Multi-provider AI widget — Gemini (free), Anthropic, or rule-based (no key)."""
    import streamlit.components.v1 as components
    escaped = json.dumps(prompt)
    html = (
        "<style>"
        "body{margin:0;background:transparent;font-family:'Inter',sans-serif;}"
        ".row{display:flex;gap:8px;margin-bottom:8px;align-items:center;}"
        "select,input{background:#0f1220;border:0.5px solid #1c2e4a;border-radius:8px;"
        "padding:9px 12px;color:#dde1e7;font-size:12px;font-family:'JetBrains Mono',monospace;outline:none;}"
        "select{cursor:pointer;}"
        "#api-key{flex:1;}"
        "#api-key::placeholder{color:#2d3650;}"
        "#run-btn{background:#10b981;color:#022c22;border:none;border-radius:8px;"
        "padding:9px 18px;font-size:12px;font-weight:700;cursor:pointer;"
        "font-family:'JetBrains Mono',monospace;white-space:nowrap;flex-shrink:0;}"
        "#run-btn:hover{background:#059669;}"
        "#run-btn:disabled{background:#1c2e4a;color:#3d4a6b;cursor:not-allowed;}"
        "#status{font-size:11px;color:#4b5780;font-family:'JetBrains Mono',monospace;min-height:16px;}"
        "#status.err{color:#f87171;} #status.ok{color:#10b981;}"
        "#result-box{display:none;background:#0a0c10;border:0.5px solid #1c2e4a;border-radius:8px;"
        "padding:12px 14px;font-size:11px;color:#8892a4;font-family:'JetBrains Mono',monospace;"
        "line-height:1.6;max-height:200px;overflow-y:auto;white-space:pre-wrap;margin-top:8px;}"
        ".hint{font-size:10px;color:#2d3650;margin-bottom:8px;line-height:1.6;}"
        "</style>"
        "<div>"
        "<div class='hint'>"
        "<b style='color:#4b5780'>Gemini is free</b> — get a key at aistudio.google.com (no credit card). "
        "No key? Choose <b style='color:#4b5780'>Rule-based</b> — same report, no AI needed."
        "</div>"
        "<div class='row'>"
        "<select id='provider' onchange='updateUI()'>"
        "<option value='gemini'>Gemini (free — AIza...)</option>"
        "<option value='anthropic'>Anthropic (Claude)</option>"
        "<option value='none'>Rule-based (no key needed)</option>"
        "</select>"
        "<input id='api-key' type='password' placeholder='Paste Gemini API key (AIza...)' autocomplete='off' />"
        "<button id='run-btn' onclick='runAI()'>Generate</button>"
        "</div>"
        "<div id='status'>Select a provider and click Generate.</div>"
        "<div id='result-box'></div>"
        "</div>"
        "<script>"
        f"const PROMPT={escaped};"
        "function updateUI(){"
        "  const p=document.getElementById('provider').value;"
        "  const ki=document.getElementById('api-key');"
        "  if(p==='none'){ki.style.display='none';}"
        "  else{ki.style.display='block';"
        "    ki.placeholder=p==='gemini'?'Paste Gemini key (AIza...)':'Paste Anthropic key (sk-ant-...)';}"
        "}"
        "async function runAI(){"
        "  const provider=document.getElementById('provider').value;"
        "  const key=document.getElementById('api-key').value.trim();"
        "  if(provider!=='none'&&!key){setStatus('Please paste your API key first.','err');return;}"
        "  const btn=document.getElementById('run-btn');"
        "  btn.disabled=true; btn.textContent='Calling…';"
        "  setStatus('Sending request…','');"
        "  try{"
        "    if(provider==='none'){"
        "      showResult(JSON.stringify({__rule_based__:true}));"
        "      setStatus('Done — click Load below to apply rule-based analysis.','ok');"
        "      btn.disabled=false; btn.textContent='Generate'; return;"
        "    }"
        "    let raw='';"
        "    if(provider==='gemini'){"
        "      const url=`https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key=${key}`;"
        "      const res=await fetch(url,{method:'POST',headers:{'Content-Type':'application/json'},"
        "        body:JSON.stringify({contents:[{parts:[{text:PROMPT}]}]})});"
        "      if(!res.ok){const e=await res.json().catch(()=>({}));throw new Error(e?.error?.message||`HTTP ${res.status}`);}"
        "      const d=await res.json();"
        "      raw=d?.candidates?.[0]?.content?.parts?.[0]?.text||'';"
        "    } else {"
        "      const res=await fetch('https://api.anthropic.com/v1/messages',{"
        "        method:'POST',"
        "        headers:{'Content-Type':'application/json','x-api-key':key,"
        "          'anthropic-version':'2023-06-01','anthropic-dangerous-direct-browser-access':'true'},"
        "        body:JSON.stringify({model:'claude-sonnet-4-20250514',max_tokens:1200,messages:[{role:'user',content:PROMPT}]}),"
        "      });"
        "      if(!res.ok){const e=await res.json().catch(()=>({}));throw new Error(e?.error?.message||`HTTP ${res.status}`);}"
        "      const d=await res.json();"
        "      raw=(d.content||[]).filter(b=>b.type==='text').map(b=>b.text).join('');"
        "    }"
        "    let cleaned=raw.trim();"
        "    if(cleaned.startsWith('```')){cleaned=cleaned.split('```')[1]||'';if(cleaned.startsWith('json'))cleaned=cleaned.slice(4);}"
        "    const parsed=JSON.parse(cleaned.trim());"
        "    setStatus('Done — copy the JSON below, paste it in the field, then click Load.','ok');"
        "    showResult(JSON.stringify(parsed,null,2));"
        "  } catch(e){"
        "    setStatus('Error: '+e.message,'err');"
        "  } finally{"
        "    btn.disabled=false; btn.textContent='Generate';"
        "  }"
        "}"
        "function showResult(text){const box=document.getElementById('result-box');box.style.display='block';box.textContent=text;}"
        "function setStatus(msg,cls){const el=document.getElementById('status');el.textContent=msg;el.className=cls;}"
        "</script>"
    )
    components.html(html, height=225, scrolling=False)


# ══════════════════════════════════════════════════════════════════════════════
# RENDER HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def hero(col, label, value, sub, css):
    col.markdown(
        f'<div class="hero-card {css}"><div class="hero-label">{label}</div>'
        f'<div class="hero-value {css}">{value}</div><div class="hero-sub">{sub}</div></div>',
        unsafe_allow_html=True,
    )


def banner(css, icon, title, body):
    st.markdown(
        f'<div class="action-banner {css}"><div class="action-icon">{icon}</div>'
        f'<div><div class="action-title {css}">{title}</div>'
        f'<div class="action-body">{body}</div></div></div>',
        unsafe_allow_html=True,
    )


# ══════════════════════════════════════════════════════════════════════════════
# EXPORT GENERATORS  — Word / PowerPoint / Excel
# ══════════════════════════════════════════════════════════════════════════════

def _chart_image_bytes(fig) -> bytes:
    """Convert a Plotly figure to PNG bytes (requires kaleido)."""
    try:
        return fig.to_image(format="png", width=900, height=320, scale=2)
    except Exception:
        return None


def build_docx(enriched, oc, guard, ai, ts, figures: dict) -> bytes:
    """Generate a professional Word report with charts embedded as images."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    import io

    doc = Document()

    # ── Page margins ──
    for section in doc.sections:
        section.top_margin    = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    # ── Styles helpers ──
    def _set_heading(para, text, size, bold=True, color=None):
        run = para.add_run(text)
        run.bold = bold
        run.font.size = Pt(size)
        if color:
            run.font.color.rgb = RGBColor(*color)
        para.paragraph_format.space_before = Pt(14)
        para.paragraph_format.space_after  = Pt(4)

    def _add_section(label: str):
        p = doc.add_paragraph()
        _set_heading(p, label.upper(), 9, color=(120, 130, 160))
        p.paragraph_format.space_before = Pt(18)
        # thin rule
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        bottom = OxmlElement('w:bottom')
        bottom.set(qn('w:val'), 'single')
        bottom.set(qn('w:sz'), '4')
        bottom.set(qn('w:space'), '2')
        bottom.set(qn('w:color'), '1E2230')
        pBdr.append(bottom)
        pPr.append(pBdr)

    def _body(text: str, bold_parts: list = None):
        p = doc.add_paragraph(style="Normal")
        p.paragraph_format.space_after = Pt(6)
        if bold_parts:
            parts = text.split("**")
            for i, part in enumerate(parts):
                run = p.add_run(part)
                run.font.size = Pt(11)
                if i % 2 == 1:
                    run.bold = True
        else:
            run = p.add_run(text)
            run.font.size = Pt(11)
        return p

    winner = enriched[enriched["is_winner"]].iloc[0]

    # ── Cover ──
    doc.add_paragraph()
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    tr = title_p.add_run("Revenue Optimization Engine")
    tr.bold = True; tr.font.size = Pt(28)
    tr.font.color.rgb = RGBColor(16, 185, 129)

    sub_p = doc.add_paragraph()
    sr = sub_p.add_run("A/B Test Management Report  ·  Growth Team")
    sr.font.size = Pt(13); sr.font.color.rgb = RGBColor(100, 116, 139)

    doc.add_paragraph()
    meta_p = doc.add_paragraph()
    mr = meta_p.add_run(f"Generated: {ts}")
    mr.font.size = Pt(10); mr.font.color.rgb = RGBColor(100, 116, 139)
    doc.add_page_break()

    # ── Executive Summary ──
    _add_section("Executive Summary")
    if ai and ai.get("executive_summary"):
        _body(ai.get("executive_summary",""))
        v = ai.get("verdict","—")
        vr = ai.get("verdict_rationale","")
        vp = doc.add_paragraph()
        vp.paragraph_format.space_after = Pt(6)
        label_run = vp.add_run(f"Verdict: {v}  —  ")
        label_run.bold = True; label_run.font.size = Pt(11)
        label_run.font.color.rgb = (RGBColor(16,185,129) if v=="SCALE" else
                                    RGBColor(245,158,11) if v=="ITERATE" else
                                    RGBColor(239,68,68))
        vp.add_run(vr).font.size = Pt(11)
    else:
        _body(f"Variant {winner['Variant']} achieves the highest revenue (${winner['Revenue_$']:,.0f}) "
              f"with the best conversion rate ({winner['Conversion_%']}%) and an RPI of ${winner['RPI']:.4f}. "
              f"The winning copy carries a {winner['brand_risk']} brand risk rating based on qualitative feedback.")

    # ── Key Metrics table ──
    _add_section("Key Metrics at a Glance")
    metrics = [
        ("Financial winner", f"Variant {winner['Variant']}"),
        ("Top RPI", f"${winner['RPI']:.4f}"),
        ("Revenue lift vs A", f"+{((winner['Revenue_$']-enriched.iloc[0]['Revenue_$'])/enriched.iloc[0]['Revenue_$']*100):.1f}%"),
        ("Annual $ at stake", f"${oc['annual_lost']:,.0f}"),
        ("Incremental / month", f"${oc['incremental_monthly']:,.0f}"),
        ("Brand risk", winner['brand_risk'].capitalize()),
        ("Data health", "WARNING" if not guard["tracking_ok"] else "Clean"),
    ]
    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = "Table Grid"
    hdr = tbl.rows[0].cells
    hdr[0].text = "Metric"; hdr[1].text = "Value"
    for cell in hdr:
        for para in cell.paragraphs:
            for run in para.runs:
                run.bold = True; run.font.size = Pt(10)
        cell._tc.get_or_add_tcPr().append(
            OxmlElement('w:shd') if False else OxmlElement('w:shd'))
    for label, value in metrics:
        row = tbl.add_row().cells
        row[0].text = label; row[1].text = value
        for para in row[0].paragraphs + row[1].paragraphs:
            for run in para.runs: run.font.size = Pt(10)
    tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
    doc.add_paragraph()

    # ── Charts ──
    _add_section("Performance Charts")
    for fname, fig in figures.items():
        img_bytes = _chart_image_bytes(fig)
        if img_bytes:
            doc.add_paragraph(fname.replace("_", " ").title())
            doc.add_picture(io.BytesIO(img_bytes), width=Inches(5.8))
            doc.add_paragraph()

    # ── Variant detail ──
    _add_section("Variant Breakdown")
    vtbl = doc.add_table(rows=1, cols=6)
    vtbl.style = "Table Grid"
    headers = ["Variant","Copy","Impressions","Conversion %","Revenue $","RPI"]
    for i, h in enumerate(headers):
        cell = vtbl.rows[0].cells[i]
        cell.text = h
        for run in cell.paragraphs[0].runs:
            run.bold = True; run.font.size = Pt(9)
    for _, row in enriched.iterrows():
        r = vtbl.add_row().cells
        r[0].text = row["Variant"]
        r[1].text = str(row["Copy"])[:40]
        r[2].text = f"{row['Impressions']:,}"
        r[3].text = f"{row['Conversion_%']}%"
        r[4].text = f"${row['Revenue_$']:,.0f}"
        r[5].text = f"${row['RPI']:.4f}"
        for c in r:
            for p in c.paragraphs:
                for run in p.runs: run.font.size = Pt(9)

    doc.add_paragraph()
    vtbl.alignment = WD_TABLE_ALIGNMENT.LEFT

    # ── Opportunity cost ──
    _add_section("Revenue Opportunity Cost")
    _body(f"By sending traffic to underperforming variants, the experiment left "
          f"**${oc['lost_revenue']:,.0f}** on the table during this period. "
          f"Scaling Variant {oc['winner_variant']} to 100% of traffic would generate "
          f"**${oc['incremental_monthly']:,.0f}/month** — equivalent to **${oc['annual_lost']:,.0f}/year** "
          f"with zero additional acquisition spend.", bold_parts=True)

    # ── AI Hypotheses ──
    if ai and ai.get("hypotheses"):
        _add_section("AI-Generated Hypotheses for Next A/B Test")
        for h in ai["hypotheses"]:
            hp = doc.add_paragraph(style="Normal")
            hp.paragraph_format.space_after = Pt(4)
            hp.add_run(f"{h['id']}: ").bold = True
            hp.runs[-1].font.size = Pt(11)
            hp.add_run(f'"{h["copy"]}"').font.size = Pt(11)
            dp = doc.add_paragraph(style="Normal")
            dp.paragraph_format.left_indent = Inches(0.3)
            dp.paragraph_format.space_after = Pt(8)
            dp.add_run(f"Rationale: {h['rationale']}\n").font.size = Pt(10)
            dp.add_run(f"Expected: {h['predicted_outcome']} · Risk: {h['risk_level']}").font.size = Pt(10)

    # ── Priority Actions ──
    if ai and ai.get("priority_actions"):
        _add_section("Priority Actions")
        for i, action in enumerate(ai["priority_actions"], 1):
            p = doc.add_paragraph(style="Normal")
            p.paragraph_format.space_after = Pt(5)
            p.paragraph_format.left_indent = Inches(0.2)
            run = p.add_run(f"{i:02d}.  {action}")
            run.font.size = Pt(11)

    # ── Data Quality ──
    _add_section("Data Quality Audit")
    if guard["missing_ctr_variants"]:
        _body(f"WARNING: CTR tracking missing for Variant(s) {', '.join(guard['missing_ctr_variants'])}. "
              f"Fix the tracking event on the desktop upsell modal before next iteration.")
    else:
        _body("All tracking signals clean. CTR data complete across all variants.")

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def build_xlsx(enriched, oc, guard, ai, ts) -> bytes:
    """Generate a multi-sheet Excel workbook with data, metrics, and hypotheses."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, Reference, LineChart
    from openpyxl.chart.series import SeriesLabel
    import io

    wb = openpyxl.Workbook()

    # ── Color palette ──
    GREEN  = "10B981"; DARK_BG = "0F1220"; MID_BG = "161B2C"
    AMBER  = "F59E0B"; RED     = "EF4444"; BLUE   = "3B82F6"
    WHITE  = "F0F4FA"; GRAY    = "8892A4"

    def _hdr_style(ws, row, cols, bg=DARK_BG, fg=WHITE, bold=True, size=11):
        fill = PatternFill("solid", fgColor=bg)
        font = Font(bold=bold, color=fg, size=size, name="Calibri")
        for col in range(1, cols + 1):
            cell = ws.cell(row=row, column=col)
            cell.fill = fill
            cell.font = font
            cell.alignment = Alignment(horizontal="center", vertical="center")

    def _cell(ws, row, col, value, bold=False, color="000000", bg=None,
              num_fmt=None, align="left"):
        c = ws.cell(row=row, column=col, value=value)
        c.font = Font(bold=bold, color=color, name="Calibri", size=10)
        c.alignment = Alignment(horizontal=align, vertical="center", wrap_text=True)
        if bg:
            c.fill = PatternFill("solid", fgColor=bg)
        if num_fmt:
            c.number_format = num_fmt
        return c

    thin = Side(style="thin", color="1E2230")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    # ════════════════════════════════════
    # SHEET 1 — Dashboard
    # ════════════════════════════════════
    ws1 = wb.active
    ws1.title = "Dashboard"
    ws1.sheet_view.showGridLines = False
    ws1.column_dimensions["A"].width = 28
    ws1.column_dimensions["B"].width = 20
    ws1.column_dimensions["C"].width = 28
    ws1.column_dimensions["D"].width = 20

    winner = enriched[enriched["is_winner"]].iloc[0]
    lift = (winner["Revenue_$"] - enriched.iloc[0]["Revenue_$"]) / enriched.iloc[0]["Revenue_$"] * 100

    # Title
    ws1["A1"] = "REVENUE OPTIMIZATION ENGINE"
    ws1["A1"].font = Font(bold=True, size=16, color=GREEN, name="Calibri")
    ws1["A2"] = f"A/B Test Management Report  ·  Generated: {ts}"
    ws1["A2"].font = Font(size=10, color=GRAY, name="Calibri")
    ws1.row_dimensions[1].height = 30
    ws1.row_dimensions[2].height = 18
    ws1.row_dimensions[3].height = 12

    # KPI cards (row 4+)
    kpis = [
        ("Financial Winner",       f"Variant {winner['Variant']}",    GREEN,  "A"),
        ("Top RPI",                f"${winner['RPI']:.4f}",           GREEN,  "C"),
        ("Revenue Lift vs A",      f"+{lift:.1f}%",                   BLUE,   "A"),
        ("Annual $ at Stake",      f"${oc['annual_lost']:,.0f}",      AMBER,  "C"),
        ("Incremental / Month",    f"${oc['incremental_monthly']:,.0f}", GREEN,"A"),
        ("Brand Risk",             winner['brand_risk'].capitalize(),  AMBER,  "C"),
    ]
    row = 4
    for i, (label, value, color, col) in enumerate(kpis):
        r = row + (i // 2) * 4
        c = 1 if col == "A" else 3
        ws1.cell(r,   c, label).font = Font(size=9, color=GRAY, bold=False, name="Calibri")
        ws1.cell(r+1, c, value).font = Font(size=15, color=color, bold=True, name="Calibri")
        ws1.cell(r+1, c).fill = PatternFill("solid", fgColor=MID_BG)
        ws1.row_dimensions[r+1].height = 24

    # ════════════════════════════════════
    # SHEET 2 — Variant Data
    # ════════════════════════════════════
    ws2 = wb.create_sheet("Variant Data")
    ws2.sheet_view.showGridLines = False
    cols_w = [12, 35, 16, 10, 14, 14, 10, 25, 10]
    col_names = ["Variant","Copy","Impressions","CTR %","Conversion %","Revenue $","RPI","Notes","Brand Risk"]
    for i, w in enumerate(cols_w, 1):
        ws2.column_dimensions[get_column_letter(i)].width = w

    ws2["A1"] = "VARIANT BREAKDOWN"
    ws2["A1"].font = Font(bold=True, size=13, color=GREEN, name="Calibri")
    ws2.row_dimensions[1].height = 24

    _hdr_style(ws2, 3, len(col_names), bg=DARK_BG, fg=WHITE)
    for i, name in enumerate(col_names, 1):
        ws2.cell(3, i, name)

    for row_i, (_, row) in enumerate(enriched.iterrows(), 4):
        bg = "E8FFF5" if row["is_winner"] else None
        fg = "065F46" if row["is_winner"] else "111827"
        _cell(ws2, row_i, 1, row["Variant"], bold=row["is_winner"], color=fg, bg=bg, align="center")
        _cell(ws2, row_i, 2, row["Copy"], color=fg, bg=bg)
        _cell(ws2, row_i, 3, row["Impressions"], color=fg, bg=bg, num_fmt="#,##0", align="right")
        ctr_val = "MISSING" if row["ctr_missing"] else row["CTR_%"]
        _cell(ws2, row_i, 4, ctr_val, color="B91C1C" if row["ctr_missing"] else fg, bg=bg, align="center")
        _cell(ws2, row_i, 5, row["Conversion_%"], color=fg, bg=bg, num_fmt='0.00"%"', align="right")
        _cell(ws2, row_i, 6, row["Revenue_$"], color=fg, bg=bg, num_fmt='"$"#,##0', align="right")
        _cell(ws2, row_i, 7, round(row["RPI"], 4), color=fg, bg=bg, num_fmt='"$"0.0000', align="right")
        _cell(ws2, row_i, 8, row["Notes"], color=fg, bg=bg)
        _cell(ws2, row_i, 9, row["brand_risk"].capitalize(), color=fg, bg=bg, align="center")
        ws2.row_dimensions[row_i].height = 20

    # Bar chart — Revenue
    chart = BarChart()
    chart.type = "col"; chart.title = "Revenue by Variant ($)"
    chart.grouping = "clustered"; chart.width = 15; chart.height = 10
    data_ref   = Reference(ws2, min_col=6, max_col=6, min_row=3, max_row=3+len(enriched))
    cats_ref   = Reference(ws2, min_col=1, max_col=1, min_row=4, max_row=3+len(enriched))
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)
    chart.series[0].graphicalProperties.solidFill = GREEN
    ws2.add_chart(chart, "A" + str(5 + len(enriched)))

    # ════════════════════════════════════
    # SHEET 3 — Opportunity Cost
    # ════════════════════════════════════
    ws3 = wb.create_sheet("Opportunity Cost")
    ws3.sheet_view.showGridLines = False
    ws3.column_dimensions["A"].width = 35
    ws3.column_dimensions["B"].width = 20

    ws3["A1"] = "OPPORTUNITY COST ANALYSIS"
    ws3["A1"].font = Font(bold=True, size=13, color=GREEN, name="Calibri")
    ws3.row_dimensions[1].height = 24

    oc_rows = [
        ("Actual Revenue (all variants)",          oc["actual_revenue"],      '"$"#,##0'),
        ("Potential Revenue (winner only)",         oc["potential_revenue"],   '"$"#,##0'),
        ("Revenue Lost This Period",                oc["lost_revenue"],        '"$"#,##0'),
        ("Daily Revenue Loss",                      oc["daily_lost"],          '"$"#,##0'),
        ("Monthly Revenue Loss",                    oc["monthly_lost"],        '"$"#,##0'),
        ("Annual Revenue at Stake",                 oc["annual_lost"],         '"$"#,##0'),
        ("Incremental Monthly (winner scaled)",     oc["incremental_monthly"], '"$"#,##0'),
        ("Annual Uplift (winner scaled)",           oc["annual_lost"],         '"$"#,##0'),
    ]
    _hdr_style(ws3, 3, 2, bg=DARK_BG, fg=WHITE)
    ws3.cell(3, 1, "Metric"); ws3.cell(3, 2, "Value")
    for i, (label, value, fmt) in enumerate(oc_rows, 4):
        ws3.row_dimensions[i].height = 20
        _cell(ws3, i, 1, label)
        _cell(ws3, i, 2, value, num_fmt=fmt, align="right",
              color=RED if "Lost" in label or "Loss" in label else GREEN)

    # 12-month projection line chart
    months = list(range(1, 13))
    ws3["A14"] = "12-Month Revenue Projection"
    ws3["A14"].font = Font(bold=True, size=11, color=WHITE, name="Calibri")
    _hdr_style(ws3, 15, 3, bg=DARK_BG, fg=WHITE)
    ws3.cell(15,1,"Month"); ws3.cell(15,2,"Current Mix"); ws3.cell(15,3,"Winner Scaled")
    for i, m in enumerate(months, 16):
        ws3.cell(i,1, m)
        ws3.cell(i,2, round(oc["baseline_monthly"] * m))
        ws3.cell(i,3, round(oc["if_scaled_monthly"] * m))
        ws3.row_dimensions[i].height = 18
    line = LineChart()
    line.title = "Cumulative Revenue: Current vs Winner Scaled"
    line.width = 18; line.height = 12
    line.add_data(Reference(ws3, min_col=2, max_col=3, min_row=15, max_row=27), titles_from_data=True)
    line.set_categories(Reference(ws3, min_col=1, min_row=16, max_row=27))
    ws3.add_chart(line, "A28")

    # ════════════════════════════════════
    # SHEET 4 — AI Hypotheses
    # ════════════════════════════════════
    ws4 = wb.create_sheet("AI Hypotheses")
    ws4.sheet_view.showGridLines = False
    ws4.column_dimensions["A"].width = 8
    ws4.column_dimensions["B"].width = 40
    ws4.column_dimensions["C"].width = 35
    ws4.column_dimensions["D"].width = 35
    ws4.column_dimensions["E"].width = 12

    ws4["A1"] = "AI-GENERATED HYPOTHESES"
    ws4["A1"].font = Font(bold=True, size=13, color=GREEN, name="Calibri")

    if ai and ai.get("executive_summary"):
        ws4["A2"] = ai.get("executive_summary","")
        ws4["A2"].font = Font(size=10, color=GRAY, name="Calibri")
        ws4["A2"].alignment = Alignment(wrap_text=True)
        ws4.row_dimensions[2].height = 48

    _hdr_style(ws4, 4, 5, bg=DARK_BG, fg=WHITE)
    for i, h in enumerate(["ID","Copy","Rationale","Expected Outcome","Risk"], 1):
        ws4.cell(4, i, h)

    if ai and ai.get("hypotheses"):
        for row_i, h in enumerate(ai["hypotheses"], 5):
            risk = h.get("risk_level","Low")
            risk_color = "065F46" if risk=="Low" else "92400E" if risk=="Medium" else "991B1B"
            risk_bg    = "D1FAE5" if risk=="Low" else "FEF3C7" if risk=="Medium" else "FEE2E2"
            _cell(ws4, row_i, 1, h.get("id",""), bold=True, color="7C3AED", align="center")
            _cell(ws4, row_i, 2, f'"{h.get("copy","")}"', bold=True)
            _cell(ws4, row_i, 3, h.get("rationale",""))
            _cell(ws4, row_i, 4, h.get("predicted_outcome",""))
            _cell(ws4, row_i, 5, risk, bold=True, color=risk_color, bg=risk_bg, align="center")
            ws4.row_dimensions[row_i].height = 40

    if ai and ai.get("priority_actions"):
        start = 5 + (len(ai.get("hypotheses",[])) or 0) + 2
        ws4.cell(start, 1, "PRIORITY ACTIONS").font = Font(bold=True, size=11, color=GREEN, name="Calibri")
        for i, action in enumerate(ai["priority_actions"], start + 1):
            ws4.cell(i, 1, f"{i - start:02d}.")
            ws4.cell(i, 2, action).font = Font(size=10, name="Calibri")
            ws4.row_dimensions[i].height = 20

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


def build_pptx(enriched, oc, guard, ai, ts, figures: dict) -> bytes:
    """Generate a professional dark-theme PowerPoint deck."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    import io

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK   = RGBColor(0x08, 0x0a, 0x0f)
    CARD   = RGBColor(0x0f, 0x12, 0x20)
    GREEN  = RGBColor(0x10, 0xb9, 0x81)
    AMBER  = RGBColor(0xf5, 0x9e, 0x0b)
    RED    = RGBColor(0xef, 0x44, 0x44)
    BLUE   = RGBColor(0x3b, 0x82, 0xf6)
    WHITE  = RGBColor(0xf0, 0xf4, 0xfa)
    MUTED  = RGBColor(0x4b, 0x57, 0x80)
    PURPLE = RGBColor(0x8b, 0x5c, 0xf6)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def _bg(slide, color=DARK):
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def _box(slide, l, t, w, h, color=CARD, radius=None):
        shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.color.rgb = RGBColor(0x1c, 0x20, 0x35)
        shp.line.width = Pt(0.5)
        return shp

    def _txt(slide, text, l, t, w, h, size=12, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf  = txb.text_frame; tf.word_wrap = wrap
        p   = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
        return txb

    def _accent_line(slide, l, t, w, color=GREEN, thickness=Pt(3)):
        line = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Emu(int(thickness)))
        line.fill.solid(); line.fill.fore_color.rgb = color
        line.line.fill.background()

    winner = enriched[enriched["is_winner"]].iloc[0]
    lift   = (winner["Revenue_$"] - enriched.iloc[0]["Revenue_$"]) / enriched.iloc[0]["Revenue_$"] * 100
    verdict_color = {"SCALE": GREEN, "ITERATE": AMBER, "KILL": RED}

    # ── Slide 1: Cover ──────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _bg(sl)
    _accent_line(sl, 0.6, 1.8, 3.0, GREEN)
    _txt(sl, "REVENUE OPTIMIZATION ENGINE", 0.6, 2.0, 9.0, 0.8, size=32, bold=True, color=GREEN)
    _txt(sl, "A/B Test Intelligence Report  ·  Growth Team", 0.6, 2.85, 9.0, 0.5, size=15, color=MUTED)
    _txt(sl, f"Generated: {ts}", 0.6, 6.9, 5.0, 0.4, size=10, color=MUTED)

    # ── Slide 2: Key Metrics ────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _bg(sl)
    _txt(sl, "KEY METRICS", 0.5, 0.25, 10, 0.45, size=10, color=MUTED)
    _txt(sl, "Performance at a Glance", 0.5, 0.6, 10, 0.7, size=24, bold=True, color=WHITE)

    kpi_data = [
        ("FINANCIAL WINNER",  f"Variant {winner['Variant']}", GREEN),
        ("TOP RPI",           f"${winner['RPI']:.4f}",        GREEN),
        ("REVENUE LIFT VS A", f"+{lift:.1f}%",                BLUE),
        ("ANNUAL $ AT STAKE", f"${oc['annual_lost']:,.0f}",   AMBER),
        ("INCR. / MONTH",     f"${oc['incremental_monthly']:,.0f}", GREEN),
        ("BRAND RISK",        winner["brand_risk"].capitalize(), AMBER),
    ]
    cols = 3; card_w = 3.8; card_h = 1.55; gap = 0.18; start_x = 0.5; start_y = 1.55
    for i, (label, value, color) in enumerate(kpi_data):
        row_i = i // cols; col_i = i % cols
        x = start_x + col_i * (card_w + gap)
        y = start_y + row_i * (card_h + gap)
        _box(sl, x, y, card_w, card_h)
        _accent_line(sl, x, y, card_w, color, Pt(3))
        _txt(sl, label, x+0.18, y+0.18, card_w-0.3, 0.35, size=8, color=MUTED)
        _txt(sl, value, x+0.18, y+0.55, card_w-0.3, 0.75, size=22, bold=True, color=color)

    # ── Slide 3: Verdict ────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _bg(sl)
    _txt(sl, "AI VERDICT", 0.5, 0.25, 10, 0.45, size=10, color=MUTED)
    verdict  = ai.get("verdict","ITERATE") if ai else decide_verdict(winner)["verdict"]
    v_color  = verdict_color.get(verdict, AMBER)
    v_body   = ai.get("verdict_rationale","") if ai else decide_verdict(winner)["body"]
    _txt(sl, verdict, 0.5, 0.65, 12, 1.2, size=56, bold=True, color=v_color)
    _box(sl, 0.5, 1.9, 12.3, 1.5)
    _txt(sl, v_body, 0.75, 2.0, 11.8, 1.3, size=14, color=WHITE, wrap=True)
    if ai and ai.get("executive_summary"):
        _txt(sl, ai["executive_summary"], 0.5, 3.65, 12.3, 1.4, size=13, color=MUTED, wrap=True)

    # ── Slide 4: Charts ─────────────────────────────────────────────────────
    for slide_title, (fig_name, fig) in zip(
        ["RPI Comparison — Normalised Performance",
         "Revenue Opportunity Cost Waterfall",
         "12-Month Revenue Projection"],
        list(figures.items())[:3]
    ):
        img_bytes = _chart_image_bytes(fig)
        if not img_bytes:
            continue
        sl = prs.slides.add_slide(blank_layout)
        _bg(sl)
        _txt(sl, "PERFORMANCE CHARTS", 0.5, 0.22, 10, 0.4, size=9, color=MUTED)
        _txt(sl, slide_title, 0.5, 0.55, 12, 0.65, size=20, bold=True, color=WHITE)
        sl.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(1.4),
                              Inches(12.3), Inches(5.6))

    # ── Slide 5: Variant table ───────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _bg(sl)
    _txt(sl, "VARIANT BREAKDOWN", 0.5, 0.22, 10, 0.4, size=9, color=MUTED)
    _txt(sl, "Full Data Table", 0.5, 0.55, 12, 0.65, size=20, bold=True, color=WHITE)

    headers  = ["VARIANT","COPY","IMPRESSIONS","CTR","CONV %","REVENUE","RPI","RISK"]
    col_widths= [0.7, 3.8, 1.3, 0.8, 0.8, 1.1, 1.0, 0.8]
    x_starts = [0.5]
    for w in col_widths[:-1]: x_starts.append(x_starts[-1] + w + 0.02)

    row_h = 0.48; start_y2 = 1.4
    for ci, (hdr, xs, cw) in enumerate(zip(headers, x_starts, col_widths)):
        _box(sl, xs, start_y2, cw, row_h, color=RGBColor(0x1c,0x20,0x35))
        _txt(sl, hdr, xs+0.06, start_y2+0.06, cw-0.1, row_h-0.1, size=8, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    for ri, (_, row) in enumerate(enriched.iterrows()):
        y = start_y2 + (ri+1) * (row_h + 0.03)
        bg_col = RGBColor(0x04,0x1f,0x13) if row["is_winner"] else CARD
        vals = [
            row["Variant"], row["Copy"][:28]+"…" if len(str(row["Copy"]))>28 else row["Copy"],
            f"{row['Impressions']:,}",
            "MISSING" if row["ctr_missing"] else f"{row['CTR_%']}%",
            f"{row['Conversion_%']}%", f"${row['Revenue_$']:,.0f}",
            f"${row['RPI']:.4f}", row["brand_risk"].capitalize(),
        ]
        for ci, (val, xs, cw) in enumerate(zip(vals, x_starts, col_widths)):
            _box(sl, xs, y, cw, row_h, color=bg_col)
            txt_color = GREEN if row["is_winner"] else (RED if val=="MISSING" else WHITE)
            _txt(sl, str(val), xs+0.05, y+0.06, cw-0.08, row_h-0.1,
                 size=8, color=txt_color, align=PP_ALIGN.CENTER)

    # ── Slide 6: AI Hypotheses ──────────────────────────────────────────────
    if ai and ai.get("hypotheses"):
        sl = prs.slides.add_slide(blank_layout)
        _bg(sl)
        _txt(sl, "AI HYPOTHESES", 0.5, 0.22, 10, 0.4, size=9, color=MUTED)
        _txt(sl, "Synthetic A/B Test Backlog — Next Iteration", 0.5, 0.55, 12, 0.65, size=20, bold=True, color=WHITE)

        hyp_colors = [PURPLE, GREEN, BLUE]
        hyp_y = 1.35
        for i, h in enumerate(ai["hypotheses"][:3]):
            risk = h.get("risk_level","Low")
            r_color = GREEN if risk=="Low" else AMBER if risk=="Medium" else RED
            _box(sl, 0.5, hyp_y, 12.3, 1.65)
            _accent_line(sl, 0.5, hyp_y, 12.3, hyp_colors[i], Pt(3))
            _txt(sl, f"{h['id']}  ·  RISK: {risk}", 0.75, hyp_y+0.12, 2.5, 0.35, size=9, color=r_color)
            _txt(sl, f'"{h.get("copy","")}"', 0.75, hyp_y+0.45, 12.0, 0.45, size=14, bold=True, color=WHITE)
            _txt(sl, h.get("rationale",""), 0.75, hyp_y+0.90, 12.0, 0.5, size=10, color=MUTED, wrap=True)
            hyp_y += 1.78

    # ── Slide 7: Priority Actions ───────────────────────────────────────────
    if ai and ai.get("priority_actions"):
        sl = prs.slides.add_slide(blank_layout)
        _bg(sl)
        _txt(sl, "ACTION PLAN", 0.5, 0.22, 10, 0.4, size=9, color=MUTED)
        _txt(sl, "Priority Actions for the Team", 0.5, 0.55, 12, 0.65, size=20, bold=True, color=WHITE)

        action_y = 1.45
        for i, action in enumerate(ai["priority_actions"][:4], 1):
            _box(sl, 0.5, action_y, 12.3, 1.1)
            _txt(sl, f"{i:02d}", 0.7, action_y+0.22, 0.6, 0.65, size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
            _txt(sl, action, 1.45, action_y+0.18, 11.1, 0.75, size=13, color=WHITE, wrap=True)
            action_y += 1.22

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def build_report(enriched, oc, guard, ai, ts) -> str:
    winner = enriched[enriched["is_winner"]].iloc[0]
    lines = [
        "REVENUE OPTIMIZATION ENGINE — MANAGEMENT REPORT",
        "=" * 54, f"Generated: {ts}", "",
        "EXECUTIVE SUMMARY", "-" * 30,
    ]
    if ai:
        lines += [ai.get("executive_summary",""),
                  f"\nVerdict: {ai.get('verdict','—')} — {ai.get('verdict_rationale','')}"]
    lines += [
        "", "WINNER", "-" * 30,
        f"Variant {winner['Variant']}: \"{winner['Copy']}\"",
        f"Revenue: ${winner['Revenue_$']:,.0f} | Conv: {winner['Conversion_%']}% | RPI: ${winner['RPI']:.4f}",
        f"Brand risk: {winner['brand_risk'].capitalize()}",
        "", "OPPORTUNITY COST", "-" * 30,
        f"Revenue lost this period: ${oc['lost_revenue']:,.0f}",
        f"Incremental monthly (if scaled): ${oc['incremental_monthly']:,.0f}",
        f"Annual at stake: ${oc['annual_lost']:,.0f}",
        "", "DATA QUALITY", "-" * 30,
        f"CTR missing: Variants {', '.join(guard['missing_ctr_variants'])} — fix tracking before next run."
        if guard["missing_ctr_variants"] else "All tracking signals clean.",
    ]
    if ai and ai.get("hypotheses"):
        lines += ["", "NEXT HYPOTHESES", "-" * 30]
        for h in ai["hypotheses"]:
            lines += [f"{h['id']}: \"{h['copy']}\" (risk: {h['risk_level']})", f"  {h['rationale']}"]
    if ai and ai.get("priority_actions"):
        lines += ["", "PRIORITY ACTIONS", "-" * 30]
        for i, a in enumerate(ai["priority_actions"], 1):
            lines.append(f"{i:02d}. {a}")
    return "\n".join(lines)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN APP
# ══════════════════════════════════════════════════════════════════════════════

ts = datetime.now().strftime("%Y-%m-%d %H:%M")

st.markdown(
    f'<div class="page-title">Revenue Optimization Engine</div>'
    f'<div class="page-sub">A/B Test Intelligence Platform &nbsp;·&nbsp; Growth Team &nbsp;·&nbsp; {ts}</div>',
    unsafe_allow_html=True,
)

tab1, tab2, tab3, tab4 = st.tabs([
    "  Input & Overview  ",
    "  Revenue Guard  ",
    "  Money Maker  ",
    "  AI Hypotheses  ",
])

# ── TAB 1 — INPUT & OVERVIEW ──────────────────────────────────────────────────
with tab1:
    st.markdown('<div class="section-head">Experiment configuration</div>', unsafe_allow_html=True)
    c1, c2, c3 = st.columns(3)
    with c1: exp_name = st.text_input("Experiment name", value="Upsell message test — desktop")
    with c2: market   = st.text_input("Market", value="USA")
    with c3: duration = st.number_input("Duration (days)", min_value=1, max_value=365, value=10)

    st.markdown('<div class="section-head">Variant data — edit directly</div>', unsafe_allow_html=True)
    st.caption("Set CTR to 0.0 if tracking failed. Add or remove rows at the bottom.")

    edited_df = st.data_editor(
        DEFAULT_DF, num_rows="dynamic", use_container_width=True,
        column_config={
            "Variant":      st.column_config.TextColumn("Variant", width=70),
            "Copy":         st.column_config.TextColumn("Copy", width=280),
            "Impressions":  st.column_config.NumberColumn("Impressions", format="%d"),
            "CTR_%":        st.column_config.NumberColumn("CTR % (0=missing)", format="%.2f%%", min_value=0),
            "Conversion_%": st.column_config.NumberColumn("Conversion %", format="%.2f%%", min_value=0),
            "Revenue_$":    st.column_config.NumberColumn("Revenue $", format="$%.0f", min_value=0),
            "Notes":        st.column_config.TextColumn("Notes / Feedback", width=220),
        },
        hide_index=True, key="data_editor",
    )

    enriched = enrich(edited_df)
    guard    = revenue_guard(enriched)
    oc       = opportunity_cost(enriched, int(duration))
    winner_r = enriched[enriched["is_winner"]].iloc[0]
    verdict  = decide_verdict(winner_r)
    lift     = (winner_r["Revenue_$"] - enriched.iloc[0]["Revenue_$"]) / enriched.iloc[0]["Revenue_$"] * 100

    st.markdown('<div class="section-head">Snapshot</div>', unsafe_allow_html=True)
    h1, h2, h3, h4, h5 = st.columns(5)
    hero(h1, "Financial winner",  f"Variant {winner_r['Variant']}", f"{winner_r['Conversion_%']}% conv · ${winner_r['Revenue_$']:,.0f}", "green")
    hero(h2, "Top RPI",           f"${winner_r['RPI']:.4f}",        "revenue per impression",    "green")
    hero(h3, "Revenue lift vs A", f"+{lift:.1f}%",                  "winner over baseline",      "blue")
    hero(h4, "Annual $ at stake", f"${oc['annual_lost']:,.0f}",     "lost to sub-optimal mix",   "amber")
    hero(h5, "Data health",
         "⚠ WARN" if not guard["tracking_ok"] else "✓ CLEAN",
         f"{guard['insight_gap_pct']:.0f}% impressions untracked" if not guard["tracking_ok"] else "All signals clean",
         "red" if not guard["tracking_ok"] else "green")

    st.markdown('<div class="section-head">Verdict</div>', unsafe_allow_html=True)
    banner(verdict["css"], verdict["icon"],
           f"ACTION REQUIRED: {verdict['verdict']} — Variant {winner_r['Variant']}",
           verdict["body"])

    st.markdown('<div class="section-head">Funnel comparison</div>', unsafe_allow_html=True)
    st.plotly_chart(chart_funnel(enriched), use_container_width=True)


# ── TAB 2 — REVENUE GUARD ─────────────────────────────────────────────────────
with tab2:
    st.markdown('<div class="section-head">Revenue guard — data quality audit</div>', unsafe_allow_html=True)

    if not guard["tracking_ok"]:
        banner("critical", "🔴",
               f"CRITICAL: CTR tracking failure — Variant(s) {', '.join(guard['missing_ctr_variants'])}",
               f"<b>{guard['insight_gap_pct']:.1f}%</b> of total impressions have no CTR data (Insight Gap). "
               "You cannot see where users drop off in the funnel. Fix the tracking event on the desktop upsell modal before next iteration. "
               "Every day without this data costs you top-of-funnel intelligence.")
    else:
        banner("success", "✅", "All tracking signals clean", "CTR data complete across all variants.")

    for alert in guard["brand_alert_variants"]:
        note = enriched.loc[enriched["Variant"] == alert["Variant"], "Notes"].values[0]
        rc   = "critical" if alert["brand_risk"] == "high" else "warning"
        banner(rc, "⚠️",
               f"BRAND RISK ({alert['brand_risk'].upper()}): Variant {alert['Variant']}",
               f"Qualitative feedback: <b>\"{note}\"</b><br>"
               "High-urgency copy drives short-term conversion but may erode long-term brand trust. "
               "Monitor support tickets and churn alongside conversion metrics.")

    st.markdown('<div class="section-head">RPI comparison</div>', unsafe_allow_html=True)
    st.plotly_chart(chart_rpi(enriched), use_container_width=True)

    st.markdown('<div class="section-head">Variant detail</div>', unsafe_allow_html=True)
    for _, row in enriched.iterrows():
        wcls  = "winner" if row["is_winner"] else ""
        wpill = (' <span style="font-size:10px;background:#022c1a;color:#10b981;padding:2px 7px;'
                 'border-radius:20px;border:0.5px solid #10b98155;font-family:JetBrains Mono;font-weight:600;margin-left:6px">WINNER</span>'
                 if row["is_winner"] else "")
        rc    = {"low":"#059669","medium":"#b45309","high":"#b91c1c"}[row["brand_risk"]]
        rb    = {"low":"#022c1a","medium":"#1a0e00","high":"#1a0404"}[row["brand_risk"]]
        ctr_h = f'<div class="vs-value">{row["CTR_%"]:.2f}%</div>' if not row["ctr_missing"] else '<div class="vs-value miss">MISSING ⚠</div>'
        bar_c = "#10b981" if row["is_winner"] else "#1c2e4a"
        st.markdown(f"""
        <div class="var-row {wcls}">
          <div class="var-name">Variant {row['Variant']}{wpill}
            <span style="font-size:10px;background:{rb};color:{rc};padding:2px 7px;border-radius:20px;
              border:0.5px solid {rc}55;font-family:JetBrains Mono;font-weight:600;margin-left:6px">
              RISK: {row['brand_risk'].upper()}
            </span>
          </div>
          <div class="var-copy">"{row['Copy']}"</div>
          <div class="var-stats">
            <div><div class="vs-label">Impressions</div><div class="vs-value">{row['Impressions']:,}</div></div>
            <div><div class="vs-label">CTR</div>{ctr_h}</div>
            <div><div class="vs-label">Conversion</div><div class="vs-value {'win' if row['is_winner'] else ''}">{row['Conversion_%']}%</div></div>
            <div><div class="vs-label">Revenue</div><div class="vs-value {'win' if row['is_winner'] else ''}">${row['Revenue_$']:,.0f}</div></div>
            <div><div class="vs-label">RPI</div><div class="vs-value {'win' if row['is_winner'] else ''}">${row['RPI']:.4f}</div></div>
            <div><div class="vs-label">Score</div><div class="vs-value">{score_variant(row)}</div></div>
          </div>
          <div class="rpi-bar-track">
            <div style="height:4px;width:{int(row['RPI_pct']*100)}%;background:{bar_c};border-radius:2px"></div>
          </div>
        </div>""", unsafe_allow_html=True)


# ── TAB 3 — MONEY MAKER ───────────────────────────────────────────────────────
with tab3:
    st.markdown('<div class="section-head">Opportunity cost — revenue lost to sub-optimal traffic split</div>', unsafe_allow_html=True)

    m1, m2, m3, m4 = st.columns(4)
    hero(m1, "Lost this period",     f"${oc['lost_revenue']:,.0f}",      f"{duration} days at current mix",  "red")
    hero(m2, "Lost / day",           f"${oc['daily_lost']:,.0f}",         "cost of not scaling winner",       "red")
    hero(m3, "Incremental / month",  f"${oc['incremental_monthly']:,.0f}", "if winner scaled to 100%",        "green")
    hero(m4, "Annual uplift",        f"${oc['annual_lost']:,.0f}",        "zero additional ad spend",         "green")

    banner("warning", "$",
           f"ACTION REQUIRED: ${oc['annual_lost']:,.0f}/year is being left on the table",
           f"<b>{oc['loser_imp']:,} impressions</b> ({oc['loser_imp']/oc['total_imp']*100:.0f}% of traffic) went to underperforming variants. "
           f"Variant {oc['winner_variant']} achieves <b>RPI ${oc['winner_rpi']:.4f}</b>. "
           f"Scaling to 100% generates <b>${oc['incremental_monthly']:,.0f}/month</b> incremental — "
           f"<b>${oc['annual_lost']:,.0f}/year</b> — with zero extra acquisition spend.")

    ca, cb = st.columns(2)
    with ca: st.plotly_chart(chart_waterfall(oc), use_container_width=True)
    with cb: st.plotly_chart(chart_projection(oc), use_container_width=True)

    st.markdown('<div class="section-head">Rollout scenario planner</div>', unsafe_allow_html=True)
    pct = st.slider("Traffic % sent to winner", 10, 100, 100, 10, format="%d%%")
    avg_loser_rpi    = enriched[~enriched["is_winner"]]["RPI"].mean()
    scenario_monthly = (oc["winner_rpi"] * oc["total_imp"] * pct/100 +
                        avg_loser_rpi * oc["loser_imp"] * (1 - pct/100)) / int(duration) * 30

    s1, s2, s3 = st.columns(3)
    hero(s1, f"Monthly at {pct}% rollout", f"${scenario_monthly:,.0f}", f"{pct}% traffic to winner", "purple")
    hero(s2, "vs current mix",             f"+${scenario_monthly - oc['baseline_monthly']:,.0f}", "incremental monthly", "blue")
    hero(s3, "Annual at this scenario",    f"${scenario_monthly*12:,.0f}", "projected full year", "green")
# ── TAB 4 — AI HYPOTHESES ────────────────────────────────────────────────────
with tab4:
    st.markdown('<div class="section-head">AI hypothesis generator — powered by AI</div>', unsafe_allow_html=True)

    st.markdown(
        f"""
        <div class="insight-card">
            The AI analyses the winning variant's qualitative feedback and generates <b>3 synthetic hypotheses</b>
            for the next test iteration. Each preserves the urgency that drove Variant {oc['winner_variant']}'s conversion
            while softening the alarmist tone to protect brand equity.
            Output is a ready-to-run test backlog grounded in data, not guesswork.
        </div>
        """,
        unsafe_allow_html=True,
    )

    # ── Session state init ───────────────────────────────────────────────────
    if "ai_result" not in st.session_state:
        st.session_state.ai_result = None
    if "ai_json_raw" not in st.session_state:
        st.session_state.ai_json_raw = ""

    # ── Prompt build + AI component ──────────────────────────────────────────
    prompt_str = build_prompt(enriched, oc, guard)
    ai_component(prompt_str)

    st.caption("📥 **Data Bridge:** Paste the generated output payload here, or simply type **DEMO** to run the local engine.")
    col_paste, col_load = st.columns([5, 1])

    with col_paste:
        # Evitamos mostrar el código JSON raro por defecto para no asustar al usuario
        default_val = st.session_state.ai_json_raw if st.session_state.ai_json_raw not in ['{"__rule_based__":true}', '{"__rule_based__": true}'] else ""

        pasted = st.text_area(
            "ai_paste", value=default_val,
            height=43, label_visibility="collapsed",
            placeholder='Paste the output data here, or type DEMO...',
            key="ai_paste_area",
        )

    with col_load:
        if st.button("Render Report", use_container_width=True):
            pasted_clean = pasted.strip()

            # El truco: Si escriben DEMO o pegan el texto original, carga sin errores
            if pasted_clean.upper() == "DEMO" or "rule_based" in pasted_clean:
                st.session_state.ai_result = make_rule_based_ai(enriched, oc, guard)
                st.session_state.ai_json_raw = "DEMO"
                st.rerun()
            else:
                try:
                    parsed = json.loads(pasted_clean)
                    st.session_state.ai_result = parsed
                    st.session_state.ai_json_raw = pasted_clean
                    st.rerun()
                except Exception:
                    st.error("⚠️ Invalid format. Paste the exact output from the console above, or type DEMO.")

    ai = st.session_state.ai_result

    # ── AI results section ───────────────────────────────────────────────────
    if ai:
        # Executive summary
        st.markdown('<div class="section-head">Executive summary</div>', unsafe_allow_html=True)
        st.markdown(
            f'<div class="insight-card">{ai.get("executive_summary", "")}<br><br>'
            f'<b>Verdict:</b> {ai.get("verdict", "—")} — {ai.get("verdict_rationale", "")}</div>',
            unsafe_allow_html=True,
        )

        # Synthetic hypotheses
        st.markdown(
            '<div class="section-head">Synthetic hypotheses for next experiments</div>',
            unsafe_allow_html=True,
        )

        risk_color_map = {
            "Low":    ("#059669", "#022c1a"),
            "Medium": ("#b45309", "#1a0e00"),
            "High":   ("#b91c1c", "#1a0404"),
        }

        for i, h in enumerate(ai.get("hypotheses", []), 1):
            risk_level = h.get("risk_level", "Low")
            risk_c, risk_b = risk_color_map.get(risk_level, risk_color_map["Low"])
            hyp_id = h.get("id", f"H{i}")

            st.markdown(
                f"""
                <div class="hyp-card">
                  <div class="hyp-num">{hyp_id} &nbsp;·&nbsp; SYNTHETIC HYPOTHESIS</div>
                  <div class="hyp-body">
                    <b>Copy:</b> "{h.get('copy', '')}"<br>
                    <b>Rationale:</b> {h.get('rationale', '')}<br>
                    <b>Expected outcome:</b> {h.get('predicted_outcome', '')}
                  </div>
                  <span style="display:inline-block;font-size:10px;background:{risk_b};color:{risk_c};
                    padding:2px 8px;border-radius:20px;border:0.5px solid {risk_c}55;margin-top:8px;
                    font-family:JetBrains Mono;font-weight:600;">
                    Risk: {risk_level}
                  </span>
                </div>
                """,
                unsafe_allow_html=True,
            )

        # Governance / system prompt transparency
        st.write("")
        with st.expander("🔍 View AI System Prompt (Governance & Brand Safety)"):
            st.markdown(
                f"""
                **System Instructions provided to the LLM:**

                > "You are an expert CRO and Growth copywriter for **Gen** (Norton, Avast, LifeLock).
                > Your task is to generate 3 new A/B test copy hypotheses based on Variant {oc['winner_variant']}'s performance.
                >
                > **CRITICAL RULE — BRAND SAFETY:** The winning variant generated negative user feedback for being 'too scary'.
                > You must soften the tone to protect our brand equity, while strictly maintaining the psychological urgency
                > that drove the conversion rate.
                > Do not use alarmist words. Focus on empowerment and digital security."
                """
            )
            st.info(
                "💡 **AI Ops Insight:** Transparent prompt architecture ensures AI outputs remain strictly aligned "
                "with Gen's brand guidelines and compliance standards before any human review."
            )

        # Priority actions (single block — no duplication)
        if ai.get("priority_actions"):
            st.markdown('<div class="section-head">Priority actions</div>', unsafe_allow_html=True)
            html_parts = ['<div class="insight-card" style="padding:16px 20px">']
            for i, action in enumerate(ai["priority_actions"], 1):
                html_parts.append(
                    f'<div style="display:flex;gap:12px;padding:9px 0;border-bottom:0.5px solid #141824;align-items:flex-start">'
                    f'<div style="width:24px;height:24px;border-radius:50%;background:#10b981;color:#022c22;'
                    f'font-size:11px;font-weight:700;display:flex;align-items:center;justify-content:center;'
                    f'flex-shrink:0;font-family:JetBrains Mono">{i:02d}</div>'
                    f'<div style="font-size:13px;color:#8892a4;line-height:1.55">{action}</div></div>'
                )
            html_parts.append("</div>")
            st.markdown("".join(html_parts), unsafe_allow_html=True)

        # Risk mitigation
        if ai.get("risk_mitigation"):
            st.markdown('<div class="section-head">Brand risk mitigation</div>', unsafe_allow_html=True)
            banner("warning", "🛡️", "Brand equity protection strategy", ai["risk_mitigation"])

    else:
        st.markdown(
            '<div class="insight-card" style="border-color:#1c2e4a;text-align:center;padding:40px 20px;">'
            '<div style="font-size:32px;margin-bottom:12px">🤖</div>'
            '<div style="font-size:14px;color:#2d3650">Press <b style="color:#8892a4">Generate</b> above to run the analysis.</div>'
            '</div>',
            unsafe_allow_html=True,
        )

    # ── Export Management Reports ────────────────────────────────────────────
    # Los botones SIEMPRE se renderizan. Si no hay AI cargada, aparecen
    # deshabilitados con tooltip explicativo en lugar de ocultarse.
    st.markdown('<div class="section-head">Export management report</div>', unsafe_allow_html=True)
    st.caption(
        "Choose a format. Each export contains the complete analysis: metrics, variant breakdown, "
        "opportunity cost, AI hypotheses, and priority actions. Word and PowerPoint embed the performance charts as images."
    )

    figs = {
        "RPI Comparison":      chart_rpi(enriched),
        "Revenue Waterfall":   chart_waterfall(oc),
        "12-Month Projection": chart_projection(oc),
        "Funnel Comparison":   chart_funnel(enriched),
    }
    date_str = ts[:10]
    ai_ready = ai is not None
    disabled_help = None if ai_ready else "Generate and load the AI analysis above to enable downloads."

    col_word, col_ppt = st.columns(2)

    with col_word:
        if ai_ready:
            try:
                docx_data = build_docx(enriched, oc, guard, ai, ts, figs)
                st.download_button(
                    label="📄 Download Word Report (.docx)",
                    data=docx_data,
                    file_name=f"revenue_report_{date_str}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    key="dl_docx",
                )
            except Exception as e:
                st.error(
                    f"⚠️ Word export failed. Run in terminal: `pip install python-docx kaleido`. Detail: {e}"
                )
        else:
            st.button(
                "📄 Download Word Report (.docx)",
                use_container_width=True,
                disabled=True,
                help=disabled_help,
                key="dl_docx_disabled",
            )

    with col_ppt:
        if ai_ready:
            try:
                pptx_data = build_pptx(enriched, oc, guard, ai, ts, figs)
                st.download_button(
                    label="📊 Download PowerPoint (.pptx)",
                    data=pptx_data,
                    file_name=f"revenue_report_{date_str}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    key="dl_pptx",
                )
            except Exception as e:
                st.error(
                    f"⚠️ PowerPoint export failed. Run in terminal: `pip install python-pptx kaleido`. Detail: {e}"
                )
        else:
            st.button(
                "📊 Download PowerPoint (.pptx)",
                use_container_width=True,
                disabled=True,
                help=disabled_help,
                key="dl_pptx_disabled",
            )
